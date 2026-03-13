# -*- coding: utf-8 -*-
from pptx import Presentation
import os

os.chdir(r"C:\Users\tsc\.openclaw\workspace\agents\参谋官")
prs = Presentation('AI数字员工汇报_总经办.pptx')
print(f'PPT共 {len(prs.slides)} 页')

# 显示每页的标题
for i, slide in enumerate(prs.slides):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) < 50 and len(text) > 0:
                title = text
                break
    print(f"第{i+1}页: {title}")
