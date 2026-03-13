# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
import os

os.chdir(r"C:\Users\tsc\.openclaw\workspace\agents\参谋官")

# 读取模板
prs = Presentation(r"h:\管理\卓越信通PPT模板.pptx")

# 读取markdown内容
with open("PPT_content.md", "r", encoding="utf-8") as f:
    content = f.read()

# 按---分割幻灯片
slides = content.split("---")
slides = [s.strip() for s in slides if s.strip()]

print(f"共 {len(slides)} 页内容要填充")

# 定义PPT内容结构
ppt_contents = [
    ("AI数字员工解决方案汇报", "——以智能采购助理为例"),
    ("目录", "1. 背景与趋势\n2. 解决方案\n3. 应用场景\n4. 价值分析"),
    ("01 背景与趋势", "传统企业管理的三大痛点\n\n• 信息滞后：依赖人工定期查询\n• 效率低下：采购员80%时间用于信息收集\n• 风险盲区：缺乏主动预警机制\n\nAI时代的答案：数字员工\n7×24小时不间断工作"),
    ("02 解决方案：OpenClaw平台", "什么是OpenClaw？\n\n• 操控浏览器、文件系统、终端\n• 集成企业IM（微信、钉钉、飞书）\n• 执行定时任务和自动化流程\n• 与企业系统无缝对接"),
    ("03 应用场景：智能采购助理", "虚拟采购助理 - 定期巡检任务\n\n• 物料停产查询：每周一自动查询\n• 价格变动监控：每日跟踪价格走势\n• 供应商资质预警：每月检查资质证照\n\n智能预警机制：\n发现异常 → 分析影响 → 推送预警 → 建议方案"),
    ("04 价值分析：提质增效", "量化收益测算\n\n• 信息及时性：滞后3-7天 → 实时/每日 ↑300%\n• 人工投入：每周8小时 → 0.5小时 ↓94%\n• 预警覆盖率：30% → 95% ↑200%\n• 应急响应：24-72小时 → <1小时 ↑2400%\n\nintangible价值：\n降低断料风险 | 释放人力 | 知识沉淀"),
]

# 清空所有幻灯片并重建
# 删除多余的幻灯片
while len(prs.slides) > len(ppt_contents):
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

# 填充内容到每页
for i, (title, body) in enumerate(ppt_contents):
    if i >= len(prs.slides):
        break
    
    slide = prs.slides[i]
    
    # 查找并更新标题和内容
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            
            # 根据位置判断是标题还是内容
            if len(text) < 30:
                # 标题框
                tf.text = title
                for p in tf.paragraphs:
                    p.font.size = Pt(32)
                    p.font.bold = True
            else:
                # 内容框
                tf.text = body
                for p in tf.paragraphs:
                    p.font.size = Pt(18)

# 保存
output_path = "AI数字员工汇报_总经办.pptx"
prs.save(output_path)
print(f"已保存: {output_path}")
print(f"共 {len(ppt_contents)} 页")
