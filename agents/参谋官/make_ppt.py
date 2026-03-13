# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

os.chdir(r"C:\Users\tsc\.openclaw\workspace\agents\参谋官")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# PPT内容
slides_content = [
    {
        "title": "AI数字员工解决方案汇报",
        "body": "——以智能采购助理为例"
    },
    {
        "title": "目录",
        "body": "01 背景与趋势\n\n02 解决方案：OpenClaw平台\n\n03 应用场景：智能采购助理\n\n04 价值分析：提质增效"
    },
    {
        "title": "01 背景与趋势",
        "body": "传统企业管理的三大痛点\n\n• 信息滞后：依赖人工定期查询供应商、物料状态\n• 效率低下：采购员80%时间用于信息收集\n• 风险盲区：缺乏主动预警机制，断料/涨价频发\n\nAI时代的答案：数字员工\n7×24小时不间断工作，严格按规则执行"
    },
    {
        "title": "02 解决方案：OpenClaw平台",
        "body": "什么是OpenClaw？\n\n一款桌面AI助手框架，让AI能够：\n\n• 操控浏览器、文件系统、终端\n• 集成企业IM（微信、钉钉、飞书）\n• 执行定时任务和自动化流程\n• 与企业系统无缝对接"
    },
    {
        "title": "03 应用场景：智能采购助理",
        "body": "虚拟采购助理 - 定期巡检任务\n\n• 物料停产查询：每周一自动查询供应商官网\n• 价格变动监控：每日跟踪重点物料价格走势\n• 供应商资质预警：每月检查资质证照有效期\n\n智能预警机制\n发现异常 → 分析影响 → 推送预警 → 建议方案"
    },
    {
        "title": "04 价值分析：提质增效",
        "body": "量化收益测算\n\n• 信息及时性：滞后3-7天 → 实时/每日  ↑300%\n• 人工投入：每周8小时 → 0.5小时     ↓94%\n• 预警覆盖率：30% → 95%            ↑200%\n• 应急响应时间：24-72小时 → <1小时  ↑2400%\n\nintangible价值\n降低断料风险 | 释放人力 | 知识沉淀 | 快速扩展"
    },
]

# 创建幻灯片
for content in slides_content:
    slide_layout = prs.slide_layouts[1]  # 标题+内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title = slide.shapes.title
    title.text = content["title"]
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # 正文
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = content["body"]
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    for p in tf.paragraphs[1:]:
        p.font.size = Pt(18)

# 保存
prs.save("AI数字员工汇报_总经办.pptx")
print("已保存: AI数字员工汇报_总经办.pptx")
print(f"共 {len(slides_content)} 页")
