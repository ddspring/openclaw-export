# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import re

os.chdir(r"C:\Users\tsc\.openclaw\workspace\agents\参谋官")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 读取完整markdown内容
with open("PPT汇报_AI数字员工.pptx.md", "r", encoding="utf-8") as f:
    content = f.read()

# 按---分割页面
pages = content.split("---")
pages = [p.strip() for p in pages if p.strip()]

print(f"共 {len(pages)} 页")

def extract_title_body(text):
    """提取标题和正文"""
    lines = text.split("\n")
    title = ""
    body_lines = []
    in_body = False
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("# ") and not title:
            title = line.replace("# ", "").strip()
            in_body = True
        elif line.startswith("## "):
            body_lines.append(("h2", line.replace("## ", "").strip()))
        elif line.startswith("### "):
            body_lines.append(("h3", line.replace("### ", "").strip()))
        elif line.startswith("| ") or line.startswith("---"):
            continue  # 跳过表格分隔符
        elif in_body:
            body_lines.append(("p", line))
    
    return title, body_lines

def add_slide(prs, title, body_lines):
    """添加一页幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题+内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # 正文
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, (line_type, text) in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        
        if line_type == "h2":
            p.font.size = Pt(24)
            p.font.bold = True
            p.space_before = Pt(12)
        elif line_type == "h3":
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_before = Pt(8)
        else:
            p.font.size = Pt(18)

# 生成每页
for i, page in enumerate(pages):
    title, body_lines = extract_title_body(page)
    if title:
        add_slide(prs, title, body_lines)
        print(f"第{len(prs.slides)}页: {title[:30]}")

# 保存
prs.save("AI数字员工汇报_总经办.pptx")
print(f"\n已保存: AI数字员工汇报_总经办.pptx")
print(f"共 {len(prs.slides)} 页")
